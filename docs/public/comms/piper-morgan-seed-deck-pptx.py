import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Create presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Define colors
BLUE_PRIMARY = RGBColor(0, 102, 204)  # #0066cc
BLUE_DARK = RGBColor(0, 68, 153)  # #004499
GRAY_DARK = RGBColor(51, 51, 51)  # #333333
GRAY_MED = RGBColor(102, 102, 102)  # #666666
TEAL = RGBColor(0, 204, 204)  # #00cccc


def add_slide_number(slide, number):
    """Add slide number to bottom right"""
    txBox = slide.shapes.add_textbox(Inches(14.5), Inches(8.3), Inches(1), Inches(0.5))
    tf = txBox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY_MED
    p.alignment = PP_ALIGN.RIGHT


# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Add blue gradient background (simplified - using solid color)
slide_1.background.fill.solid()
slide_1.background.fill.fore_color.rgb = BLUE_PRIMARY

# Logo placeholder
left = Inches(7.25)
top = Inches(1)
width = Inches(1.5)
height = Inches(1.5)
shape = slide_1.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
shape.line.color.rgb = RGBColor(255, 255, 255)

# Add logo instruction text
logo_text = slide_1.shapes.add_textbox(Inches(6.5), Inches(2.6), Inches(3), Inches(0.5))
tf = logo_text.text_frame
p = tf.paragraphs[0]
p.text = "[Replace with Piper Morgan Logo]"
p.font.size = Pt(12)
p.font.color.rgb = RGBColor(200, 200, 200)
p.alignment = PP_ALIGN.CENTER

# Title
title_box = slide_1.shapes.add_textbox(Inches(2), Inches(3.2), Inches(12), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Piper Morgan"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide_1.shapes.add_textbox(Inches(2), Inches(4.3), Inches(12), Inches(0.6))
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = "The AI Chief of Staff for Product Managers"
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Tagline
tagline_box = slide_1.shapes.add_textbox(Inches(2), Inches(5), Inches(12), Inches(0.6))
tf = tagline_box.text_frame
p = tf.paragraphs[0]
p.text = "Making every PM as effective as the best PM you've ever worked with"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(230, 230, 230)
p.alignment = PP_ALIGN.CENTER

# Credentials
cred_box = slide_1.shapes.add_textbox(Inches(2), Inches(6.5), Inches(12), Inches(0.8))
tf = cred_box.text_frame
p = tf.paragraphs[0]
p.text = "Christian Crumlish • Senior Director Product at CloudOn ($100M Exit)"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(200, 200, 200)
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "Author, Product Management for UX People"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(200, 200, 200)
p.alignment = PP_ALIGN.CENTER

add_slide_number(slide_1, 1)

# Slide 2: Problem
slide_2 = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide_2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "The Problem: PMs Are Drowning in Busywork"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = BLUE_PRIMARY

# Big stat
stat_box = slide_2.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(14), Inches(1.5)
)
stat_box.fill.solid()
stat_box.fill.fore_color.rgb = RGBColor(248, 249, 250)
stat_box.line.color.rgb = BLUE_PRIMARY
stat_box.line.width = Pt(3)

stat_text = slide_2.shapes.add_textbox(Inches(1), Inches(2.2), Inches(14), Inches(1.2))
tf = stat_text.text_frame
p = tf.paragraphs[0]
p.text = "60% of PM time spent on administrative tasks, not strategic thinking"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = GRAY_DARK
p.alignment = PP_ALIGN.CENTER

# Bullet points
bullets = slide_2.shapes.add_textbox(Inches(1), Inches(4), Inches(14), Inches(3.5))
tf = bullets.text_frame
tf.margin_left = Inches(0.5)

points = [
    "Current tools are either toys ($20/month ChatGPT wrappers) or tanks ($400/month enterprise suites)",
    "No AI truly understands PM work - they generate documents, not decisions",
    'The "missing middle": 100K+ companies with 5-50 person product teams have no good options',
    "PMs need partners that understand context, not just content",
]

for point in points:
    p = tf.add_paragraph()
    p.text = f"• {point}"
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY_DARK
    p.level = 0
    p.space_after = Pt(12)

add_slide_number(slide_2, 2)

# Slide 3: Solution
slide_3 = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide_3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "The Solution: AI That Works Like a Senior PM"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = BLUE_PRIMARY

# Three pillars
for i, (icon, title, desc) in enumerate(
    [
        (
            "🧠",
            "Spatial Intelligence",
            "Understands relationships\nlike humans do -\nnot just data points",
        ),
        ("🔌", "MCP Integration", "Connects to your entire\nPM stack in minutes,\nnot months"),
        ("📊", "Systematic Excellence", "Codified methodology,\nnot random prompts"),
    ]
):
    left = Inches(1.5 + i * 4.5)

    # Icon box
    icon_shape = slide_3.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(1), Inches(2), Inches(1.5), Inches(1.5)
    )
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = BLUE_PRIMARY

    # Title
    title_text = slide_3.shapes.add_textbox(left, Inches(3.8), Inches(3.5), Inches(0.5))
    tf = title_text.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Description
    desc_text = slide_3.shapes.add_textbox(left, Inches(4.4), Inches(3.5), Inches(1.5))
    tf = desc_text.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY_MED
    p.alignment = PP_ALIGN.CENTER

# Bottom explanation
explanation = slide_3.shapes.add_textbox(Inches(1), Inches(6.5), Inches(14), Inches(1))
tf = explanation.text_frame
p = tf.paragraphs[0]
p.text = "Piper develops environmental awareness - understanding your Slack as a navigable space,"
p.font.size = Pt(16)
p.font.color.rgb = GRAY_MED
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "tracking attention patterns, and maintaining context across conversations"
p.font.size = Pt(16)
p.font.color.rgb = GRAY_MED
p.alignment = PP_ALIGN.CENTER

add_slide_number(slide_3, 3)

# Slide 4: Market Opportunity
slide_4 = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide_4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Market Opportunity: Explosive Growth"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = BLUE_PRIMARY

# Market boxes
market_data = [
    ("$50.3B", "TAM by 2030", "45.8% CAGR"),
    ("$5B", "SAM", "Underserved mid-market"),
    ("$500M", "SOM in 5 years", "10% of SAM"),
]

for i, (number, label, sublabel) in enumerate(market_data):
    left = Inches(1.5 + i * 4.5)

    # Box
    box = slide_4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2), Inches(4), Inches(2.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(240, 249, 255)
    box.line.color.rgb = BLUE_PRIMARY
    box.line.width = Pt(2)

    # Number
    num_text = slide_4.shapes.add_textbox(left, Inches(2.3), Inches(4), Inches(0.8))
    tf = num_text.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLUE_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # Label
    label_text = slide_4.shapes.add_textbox(left, Inches(3.1), Inches(4), Inches(0.5))
    tf = label_text.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_DARK
    p.alignment = PP_ALIGN.CENTER

    # Sublabel
    sublabel_text = slide_4.shapes.add_textbox(left, Inches(3.6), Inches(4), Inches(0.5))
    tf = sublabel_text.text_frame
    p = tf.paragraphs[0]
    p.text = sublabel
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY_MED
    p.alignment = PP_ALIGN.CENTER

# Additional points
points_box = slide_4.shapes.add_textbox(Inches(1), Inches(5.2), Inches(14), Inches(2))
tf = points_box.text_frame
points = [
    "Bottom-up adoption: Individual PMs → Teams → Enterprise",
    "78% of enterprises using AI but only 15% happy with results",
    "ChatPRD proved market (17K users, unfunded) but can't scale",
]
for point in points:
    p = tf.add_paragraph()
    p.text = f"• {point}"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY_DARK

add_slide_number(slide_4, 4)

# Slide 5: Product Magic
slide_5 = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide_5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "The Magic: Multi-Agent Orchestration"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = BLUE_PRIMARY

# Three-step process
steps = [
    ("💬", "PM Describes Need", '"Help me understand why\nour onboarding\nconversion dropped"'),
    ("🤖", "Piper Orchestrates", "Multiple AI agents\nworking in parallel\nwith spatial awareness"),
    ("✅", "Complete Delivery", "GitHub issues,\nanalysis, and\nrecommendations ready"),
]

for i, (icon, title, desc) in enumerate(steps):
    left = Inches(1.5 + i * 4.5)

    # Step number/icon placeholder
    step_text = slide_5.shapes.add_textbox(left + Inches(1), Inches(2.5), Inches(2), Inches(1))
    tf = step_text.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(48)
    p.alignment = PP_ALIGN.CENTER

    # Title
    title_text = slide_5.shapes.add_textbox(left, Inches(3.8), Inches(4), Inches(0.5))
    tf = title_text.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Description
    desc_text = slide_5.shapes.add_textbox(left, Inches(4.4), Inches(4), Inches(1.2))
    tf = desc_text.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY_MED
    p.alignment = PP_ALIGN.CENTER

    # Arrow
    if i < 2:
        arrow = slide_5.shapes.add_textbox(left + Inches(4), Inches(3), Inches(0.5), Inches(0.5))
        tf = arrow.text_frame
        p = tf.paragraphs[0]
        p.text = "→"
        p.font.size = Pt(30)
        p.font.color.rgb = BLUE_PRIMARY

# Example box
example_box = slide_5.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(6.2), Inches(10), Inches(1.2)
)
example_box.fill.solid()
example_box.fill.fore_color.rgb = RGBColor(240, 249, 255)

example_text = slide_5.shapes.add_textbox(Inches(3), Inches(6.4), Inches(10), Inches(0.8))
tf = example_text.text_frame
p = tf.paragraphs[0]
p.text = "Real Example: 4 hours of PM work completed in 30 minutes"
p.font.size = Pt(18)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "From our current production usage"
p.font.size = Pt(14)
p.font.color.rgb = GRAY_MED
p.alignment = PP_ALIGN.CENTER

add_slide_number(slide_5, 5)

# Slide 6: Traction
slide_6 = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide_6.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Traction & Validation"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = BLUE_PRIMARY

# Traction metrics
metrics = [
    ("4+", "Months Systematic\nDevelopment"),
    ("635+", "PMs Following\nthe Build"),
    ("20+", "Published\nCase Studies"),
    ("100%", "Open Source\nTransparency"),
]

for i, (number, label) in enumerate(metrics):
    left = Inches(1.5 + (i % 2) * 7)
    top = Inches(2 + (i // 2) * 2.5)

    # Metric box
    box = slide_6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6), Inches(2))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(248, 249, 250)

    # Number
    num_text = slide_6.shapes.add_textbox(left, top + Inches(0.3), Inches(6), Inches(0.8))
    tf = num_text.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLUE_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # Label
    label_text = slide_6.shapes.add_textbox(left, top + Inches(1.1), Inches(6), Inches(0.8))
    tf = label_text.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_MED
    p.alignment = PP_ALIGN.CENTER

add_slide_number(slide_6, 6)

# Slide 7: Why Now
slide_7 = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide_7.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Why Now?"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = BLUE_PRIMARY

# Timing points
timing_box = slide_7.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(4.5))
tf = timing_box.text_frame

points = [
    "MCP protocol just launched (Nov 2024) - first mover advantage for seamless integrations",
    "78% of enterprises using AI but only 15% happy with results",
    "ChatPRD proved market (17K users, unfunded) but can't scale to enterprise",
    "Post-COVID remote work makes AI collaboration essential, not optional",
    "$100B+ in AI funding in 2024 - investors seeking applications, not infrastructure",
]

for point in points:
    p = tf.add_paragraph()
    p.text = f"• {point}"
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY_DARK
    p.space_after = Pt(16)

# Bottom callout
callout = slide_7.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(6.8), Inches(10), Inches(1)
)
callout.fill.solid()
callout.fill.fore_color.rgb = RGBColor(240, 249, 255)

callout_text = slide_7.shapes.add_textbox(Inches(3), Inches(7), Inches(10), Inches(0.6))
tf = callout_text.text_frame
p = tf.paragraphs[0]
p.text = "The window for defining how AI augments PM work is open NOW"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = BLUE_PRIMARY
p.alignment = PP_ALIGN.CENTER

add_slide_number(slide_7, 7)

# Slide 8: Business Model
slide_8 = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide_8.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Business Model: Open Core SaaS"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = BLUE_PRIMARY

# Pricing tiers
tiers = [
    ("Free", "Individual PM", ["Core features", "Community support", "Public repos"]),
    ("Pro", "$49/PM/month", ["Team features", "Priority support", "Advanced analytics"]),
    ("Enterprise", "$149/PM/month", ["SSO & compliance", "Custom training", "Dedicated support"]),
]

for i, (tier, price, features) in enumerate(tiers):
    left = Inches(1.5 + i * 4.5)

    # Tier box
    box = slide_8.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2), Inches(4), Inches(4)
    )
    box.fill.solid()
    if i == 1:  # Highlight Pro tier
        box.fill.fore_color.rgb = RGBColor(240, 249, 255)
        box.line.color.rgb = BLUE_PRIMARY
        box.line.width = Pt(2)
    else:
        box.fill.fore_color.rgb = RGBColor(248, 249, 250)
        box.line.color.rgb = RGBColor(220, 220, 220)

    # Tier name
    tier_text = slide_8.shapes.add_textbox(left, Inches(2.3), Inches(4), Inches(0.5))
    tf = tier_text.text_frame
    p = tf.paragraphs[0]
    p.text = tier
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_PRIMARY if i == 1 else GRAY_MED
    p.alignment = PP_ALIGN.CENTER

    # Price
    price_text = slide_8.shapes.add_textbox(left, Inches(2.9), Inches(4), Inches(0.6))
    tf = price_text.text_frame
    p = tf.paragraphs[0]
    p.text = price
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Features
    features_text = slide_8.shapes.add_textbox(left, Inches(3.8), Inches(4), Inches(1.8))
    tf = features_text.text_frame
    for feature in features:
        p = tf.add_paragraph()
        p.text = f"✓ {feature}"
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY_DARK
        p.alignment = PP_ALIGN.CENTER

# Path to $100M ARR
path_text = slide_8.shapes.add_textbox(Inches(1), Inches(6.5), Inches(14), Inches(0.6))
tf = path_text.text_frame
p = tf.paragraphs[0]
p.text = "Path to $100M ARR in 5 years"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = BLUE_PRIMARY
p.alignment = PP_ALIGN.CENTER

add_slide_number(slide_8, 8)

# Slide 9: Competition
slide_9 = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide_9.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Competition & Our Moats"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = BLUE_PRIMARY

# Competition matrix (simplified as text positioning)
matrix_box = slide_9.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(14), Inches(3.5)
)
matrix_box.fill.solid()
matrix_box.fill.fore_color.rgb = RGBColor(248, 249, 250)
matrix_box.line.color.rgb = RGBColor(200, 200, 200)

# Axis labels
axis_labels = [
    (Inches(1.2), Inches(2.2), "High AI\nSophistication"),
    (Inches(1.2), Inches(5), "Simple Tools"),
    (Inches(14.2), Inches(5), "High Price"),
    (Inches(1.2), Inches(3.5), "Low Price"),
]

for left, top, text in axis_labels:
    label = slide_9.shapes.add_textbox(left, top, Inches(1.5), Inches(0.5))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY_MED

# Competitors
competitors = [
    (Inches(3), Inches(4.5), "ChatGPT", False),
    (Inches(4), Inches(4), "ChatPRD", False),
    (Inches(11), Inches(3), "Productboard", False),
    (Inches(10), Inches(3.3), "Aha!", False),
    (Inches(7), Inches(2.5), "Piper Morgan", True),
]

for left, top, name, is_piper in competitors:
    comp = slide_9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2), Inches(0.5))
    comp.fill.solid()
    if is_piper:
        comp.fill.fore_color.rgb = BLUE_PRIMARY
    else:
        comp.fill.fore_color.rgb = RGBColor(255, 255, 255)
        comp.line.color.rgb = RGBColor(200, 200, 200)

    comp_text = slide_9.shapes.add_textbox(left, top, Inches(2), Inches(0.5))
    tf = comp_text.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(12)
    p.font.bold = is_piper
    p.font.color.rgb = RGBColor(255, 255, 255) if is_piper else GRAY_DARK
    p.alignment = PP_ALIGN.CENTER

# Moats
moats_data = [
    ("🏰", "MCP First-Mover", "Universal integration protocol"),
    ("🌐", "Open Source Community", "Network effects & trust"),
    ("🧠", "Spatial Intelligence", "Patents pending"),
]

for i, (icon, title, desc) in enumerate(moats_data):
    left = Inches(2 + i * 4.5)

    moat_text = slide_9.shapes.add_textbox(left, Inches(6), Inches(3.5), Inches(1.2))
    tf = moat_text.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY_MED
    p.alignment = PP_ALIGN.CENTER

add_slide_number(slide_9, 9)

# Slide 10: Team
slide_10 = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide_10.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Team & Advisors"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = BLUE_PRIMARY

# Founder section
founder_title = slide_10.shapes.add_textbox(Inches(1), Inches(1.8), Inches(14), Inches(0.6))
tf = founder_title.text_frame
p = tf.paragraphs[0]
p.text = "Christian Crumlish - Founder & CEO"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = GRAY_DARK

# Founder credentials
creds_box = slide_10.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(2.8))
tf = creds_box.text_frame
credentials = [
    "Senior Director of Product at CloudOn - scaled to 9M users, $100M Dropbox acquisition",
    "Delivered COVID19.CA.GOV to 40M residents under crisis timeline",
    "Proven ability to ship at scale in complex regulatory environments",
    'Author: "Product Management for UX People" (Rosenfeld Media, 2022)',
    "20+ years Silicon Valley experience (Yahoo!, AOL, 18F)",
]

for cred in credentials:
    p = tf.add_paragraph()
    p.text = f"• {cred}"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_DARK
    p.space_after = Pt(8)

# Hiring plan
hiring_title = slide_10.shapes.add_textbox(Inches(1), Inches(5.5), Inches(14), Inches(0.5))
tf = hiring_title.text_frame
p = tf.paragraphs[0]
p.text = "With Funding, We'll Hire:"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = GRAY_DARK

# Hiring boxes
hires = [
    ("CTO", "Ex-FAANG from network"),
    ("Head of AI/ML", "Spatial intelligence expert"),
    ("Engineer #1", "Full-stack senior"),
    ("Engineer #2", "Platform specialist"),
]

for i, (role, desc) in enumerate(hires):
    left = Inches(1.5 + i * 3.5)

    hire_box = slide_10.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(6.2), Inches(3), Inches(1.2)
    )
    hire_box.fill.solid()
    hire_box.fill.fore_color.rgb = RGBColor(248, 249, 250)

    hire_text = slide_10.shapes.add_textbox(left, Inches(6.4), Inches(3), Inches(0.8))
    tf = hire_text.text_frame
    p = tf.paragraphs[0]
    p.text = role
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY_MED
    p.alignment = PP_ALIGN.CENTER

add_slide_number(slide_10, 10)

# Slide 11: The Ask
slide_11 = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide_11.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "The Ask"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = BLUE_PRIMARY

# Main ask
ask_box = slide_11.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(1.8), Inches(6), Inches(1)
)
ask_box.fill.solid()
ask_box.fill.fore_color.rgb = RGBColor(240, 249, 255)
ask_box.line.color.rgb = BLUE_PRIMARY
ask_box.line.width = Pt(3)

ask_text = slide_11.shapes.add_textbox(Inches(5), Inches(2), Inches(6), Inches(0.6))
tf = ask_text.text_frame
p = tf.paragraphs[0]
p.text = "$3M Seed Round"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = BLUE_PRIMARY
p.alignment = PP_ALIGN.CENTER

# Use of funds
funds_title = slide_11.shapes.add_textbox(Inches(1), Inches(3.2), Inches(6), Inches(0.5))
tf = funds_title.text_frame
p = tf.paragraphs[0]
p.text = "Use of Funds:"
p.font.size = Pt(20)
p.font.bold = True

funds = [
    ("40%", "Core team (4 key hires)"),
    ("30%", "Product development & AI costs"),
    ("20%", "Go-to-market experiments"),
    ("10%", "Operations & legal"),
]

for i, (percent, use) in enumerate(funds):
    fund_text = slide_11.shapes.add_textbox(
        Inches(1.5), Inches(3.8 + i * 0.6), Inches(6), Inches(0.5)
    )
    tf = fund_text.text_frame
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = f"{percent} - "
    run1.font.bold = True
    run1.font.color.rgb = BLUE_PRIMARY
    run1.font.size = Pt(16)
    run2 = p.add_run()
    run2.text = use
    run2.font.size = Pt(16)

# Milestones
milestones_title = slide_11.shapes.add_textbox(Inches(8), Inches(3.2), Inches(6), Inches(0.5))
tf = milestones_title.text_frame
p = tf.paragraphs[0]
p.text = "Milestones:"
p.font.size = Pt(20)
p.font.bold = True

milestones = [
    ("6 months", "1,000 active users"),
    ("12 months", "$50K MRR"),
    ("18 months", "Series A ready"),
]

for i, (time, goal) in enumerate(milestones):
    milestone_box = slide_11.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(3.8 + i * 1.2), Inches(5), Inches(0.9)
    )
    milestone_box.fill.solid()
    milestone_box.fill.fore_color.rgb = RGBColor(248, 249, 250)

    milestone_text = slide_11.shapes.add_textbox(
        Inches(8.5), Inches(3.9 + i * 1.2), Inches(5), Inches(0.7)
    )
    tf = milestone_text.text_frame
    p = tf.paragraphs[0]
    p.text = time
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = goal
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER

add_slide_number(slide_11, 11)

# Slide 12: Vision
slide_12 = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide_12.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "The Big Vision"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = BLUE_PRIMARY

# Vision timeline
vision_items = [
    ("Year 1", "Best AI assistant\nfor PMs"),
    ("Year 3", "Operating system\nfor product teams"),
    ("Year 5", "How all digital\nproducts get built"),
]

for i, (year, vision) in enumerate(vision_items):
    left = Inches(2 + i * 4.5)

    year_text = slide_12.shapes.add_textbox(left, Inches(2.2), Inches(3), Inches(0.5))
    tf = year_text.text_frame
    p = tf.paragraphs[0]
    p.text = year
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    vision_text = slide_12.shapes.add_textbox(left, Inches(2.8), Inches(3), Inches(0.8))
    tf = vision_text.text_frame
    p = tf.paragraphs[0]
    p.text = vision
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY_MED
    p.alignment = PP_ALIGN.CENTER

# Main message
message_box = slide_12.shapes.add_textbox(Inches(2), Inches(4.5), Inches(12), Inches(2))
tf = message_box.text_frame
p = tf.paragraphs[0]
p.text = "We're not replacing PMs —"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = BLUE_PRIMARY
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "we're giving them superpowers"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = BLUE_PRIMARY
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = "AI doesn't need to think like a PM."
p.font.size = Pt(18)
p.font.color.rgb = GRAY_MED
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "It needs to work like one."
p.font.size = Pt(18)
p.font.color.rgb = GRAY_MED
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "That's about orchestration, not generation."
p.font.size = Pt(18)
p.font.color.rgb = GRAY_MED
p.alignment = PP_ALIGN.CENTER

# Contact CTA
contact_box = slide_12.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(7), Inches(10), Inches(1.2)
)
contact_box.fill.solid()
contact_box.fill.fore_color.rgb = BLUE_PRIMARY

contact_text = slide_12.shapes.add_textbox(Inches(3), Inches(7.1), Inches(10), Inches(1))
tf = contact_text.text_frame
p = tf.paragraphs[0]
p.text = "Join us in building the future of product management"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "christian@pipermorgan.ai"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(230, 230, 230)
p.alignment = PP_ALIGN.CENTER

add_slide_number(slide_12, 12)

# Save the presentation
output = io.BytesIO()
prs.save(output)
output.seek(0)

# The file is now ready to be downloaded
print("PowerPoint presentation created successfully!")
print("File size:", len(output.getvalue()), "bytes")
print("\nInstructions for use:")
print("1. Download this file as 'Piper_Morgan_Seed_Deck.pptx'")
print("2. Replace placeholder logo with your Mel Choyce-designed logo")
print("3. Apply Hoss Round font to headings if available")
print("4. Adjust colors to match your exact brand palette")
print("\nKey corrections made:")
print("- Christian's role at CloudOn properly positioned")
print("- Government experience reframed for impact")
print("- Spatial intelligence better explained")

# Return the file for download
output
